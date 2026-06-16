#!/usr/bin/env python3
"""Generate the FieldSprout investor deck as a .pptx file.

Usage: python scripts/build_investor_deck.py [output_path]
Default output: FieldSprout-Investor-Deck.pptx (repo root)
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Brand palette ----
INK      = RGBColor(0x0F, 0x17, 0x2A)   # near-black navy
PRIMARY  = RGBColor(0x4F, 0x46, 0xE5)   # indigo
PRIMARY_D= RGBColor(0x37, 0x30, 0xA3)   # deep indigo
ACCENT   = RGBColor(0x22, 0xC5, 0x5E)   # green
SLATE    = RGBColor(0x47, 0x55, 0x69)   # body text
MUTED    = RGBColor(0x94, 0xA3, 0xB8)   # muted
LIGHT    = RGBColor(0xF1, 0xF5, 0xF9)   # light panel
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

# 16:9
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.05):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        if isinstance(para, tuple):
            para = [para]
        for (txt, size, color, bold) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = "Calibri"
    return tb


def kicker(s, label, color=PRIMARY):
    bar = rect(s, Inches(0.9), Inches(0.62), Inches(0.32), Inches(0.10), ACCENT)
    text(s, Inches(1.32), Inches(0.42), Inches(8), Inches(0.5),
         [[(label.upper(), 13, color, True)]])


def header(s, title, sub=None):
    kicker(s, "FieldSprout")
    text(s, Inches(0.9), Inches(0.95), Inches(11.5), Inches(1.0),
         [[(title, 32, INK, True)]])
    if sub:
        text(s, Inches(0.92), Inches(1.75), Inches(11.5), Inches(0.6),
             [[(sub, 16, SLATE, False)]])


def footer(s, n):
    text(s, Inches(0.9), Inches(7.02), Inches(6), Inches(0.4),
         [[("FieldSprout  ·  Confidential", 9, MUTED, False)]])
    text(s, Inches(11.6), Inches(7.02), Inches(0.9), Inches(0.4),
         [[(str(n), 9, MUTED, True)]], align=PP_ALIGN.RIGHT)


# =====================================================================
# 1. TITLE
# =====================================================================
s = slide()
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, 0, Inches(0.28), SH, ACCENT)
# subtle band
rect(s, 0, Inches(4.9), SW, Inches(0.04), PRIMARY)
text(s, Inches(1.0), Inches(0.9), Inches(8), Inches(0.5),
     [[("FieldSprout", 22, WHITE, True)]])
text(s, Inches(1.0), Inches(2.5), Inches(11.2), Inches(2.2),
     [[("The AI Marketing", 52, WHITE, True)],
      [("Command Center for Local Pros", 52, ACCENT, True)]],
     line_spacing=1.0)
text(s, Inches(1.02), Inches(5.15), Inches(10.5), Inches(1.0),
     [[("One dashboard to run Google Ads, Facebook Ads, and your website — ", 18, RGBColor(0xCB,0xD5,0xE1), False)],
      [("optimized 24/7 by autonomous AI agents.", 18, RGBColor(0xCB,0xD5,0xE1), False)]])
text(s, Inches(1.02), Inches(6.5), Inches(11), Inches(0.5),
     [[("Investor Overview  ·  2026  ·  Confidential", 12, MUTED, True)]])

# =====================================================================
# 2. THE PROBLEM
# =====================================================================
s = slide()
header(s, "Local service businesses are losing the marketing game",
       "32M+ home & field service businesses in the US — most market blind.")
problems = [
    ("fa-coins", "Wasted ad spend", "20–40% of Google & Facebook budgets are burned on the wrong clicks, with no one watching daily."),
    ("fa-people-roof", "Agencies are expensive & opaque", "$1,000–$3,000/mo retainers, slow humans, and reports the owner can't understand."),
    ("fa-puzzle-piece", "Tools are fragmented", "Ads, social, and website live in separate logins — nothing talks to each other."),
    ("fa-user-clock", "Owners have no time", "Plumbers and HVAC techs run the business by day; marketing happens never."),
]
x0 = Inches(0.9); gap = Inches(0.35)
cw = Inches(2.78); ch = Inches(3.3)
for i, (ic, t, d) in enumerate(problems):
    x = x0 + i * (cw + gap)
    card = rect(s, x, Inches(2.6), cw, ch, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, Inches(2.6), cw, Inches(0.12), PRIMARY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x + Inches(0.25), Inches(2.95), cw - Inches(0.5), Inches(0.7),
         [[(t, 17, INK, True)]])
    text(s, x + Inches(0.25), Inches(3.75), cw - Inches(0.5), Inches(2.0),
         [[(d, 13, SLATE, False)]], line_spacing=1.15)
text(s, Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.6),
     [[("The result: ", 15, INK, True), ("billions in local ad spend wasted every year, and owners who give up on marketing entirely.", 15, SLATE, False)]])
footer(s, 2)

# =====================================================================
# 3. OUR SOLUTION
# =====================================================================
s = slide()
header(s, "FieldSprout is the command center that runs it all",
       "Connect your accounts in minutes. AI agents do the rest — continuously.")
sol = [
    ("Connect", "Link Google Ads, Facebook Ads, and WordPress in under 5 minutes. No setup project, no consultant."),
    ("Set your goals", "A guided wizard asks plain-English questions: target cost per lead, monthly budget, growth goal."),
    ("AI agents optimize 24/7", "Strategic, operational, and tactical agents adjust bids, budgets, keywords, audiences, and content."),
    ("You stay in control", "Plain-English digests, one-tap approvals, and guardrails you set. Cancel the agency."),
]
y = Inches(2.7)
for i, (t, d) in enumerate(sol):
    yy = y + i * Inches(0.95)
    num = rect(s, Inches(0.95), yy, Inches(0.62), Inches(0.62), PRIMARY, shape=MSO_SHAPE.OVAL)
    text(s, Inches(0.95), yy + Inches(0.04), Inches(0.62), Inches(0.55),
         [[(str(i+1), 20, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(1.85), yy - Inches(0.03), Inches(3.4), Inches(0.7),
         [[(t, 18, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(5.4), yy - Inches(0.03), Inches(7.0), Inches(0.9),
         [[(d, 14, SLATE, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
footer(s, 3)

# =====================================================================
# 4. HOW IT WORKS — AGENT ARCHITECTURE
# =====================================================================
s = slide()
header(s, "A three-tier AI agent system — built to scale across channels",
       "Each new channel plugs in under one strategic brain.")
# Strategic
rect(s, Inches(2.3), Inches(2.5), Inches(8.7), Inches(0.95), PRIMARY_D, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(2.3), Inches(2.55), Inches(8.7), Inches(0.85),
     [[("STRATEGIC  ", 16, WHITE, True), ("· the umbrella over every decision — moves budget to the best-performing channel", 13, RGBColor(0xC7,0xD2,0xFE), False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# Operational
rect(s, Inches(2.3), Inches(3.65), Inches(8.7), Inches(0.95), PRIMARY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(2.3), Inches(3.7), Inches(8.7), Inches(0.85),
     [[("OPERATIONAL  ", 16, WHITE, True), ("· channel-specific strategy for Google, Facebook, and Website", 13, RGBColor(0xE0,0xE7,0xFF), False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# Tactical channels
chans = [("Google Ads", ACCENT), ("Facebook Ads", ACCENT), ("Website / SEO", ACCENT)]
cw = Inches(2.78); x0 = Inches(2.3); gp = Inches(0.18)
for i, (nm, col) in enumerate(chans):
    x = x0 + i * (cw + gp)
    rect(s, x, Inches(4.8), cw, Inches(0.9), INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x, Inches(4.85), cw, Inches(0.8),
         [[("TACTICAL", 11, ACCENT, True)], [(nm, 14, WHITE, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
text(s, Inches(0.9), Inches(6.1), Inches(11.5), Inches(0.8),
     [[("Adaptive cadence: ", 14, INK, True),
       ("agents run more often when there's something to fix and back off when an account is steady — keeping cost and API usage low as we scale.", 14, SLATE, False)]],
     line_spacing=1.15)
footer(s, 4)

# =====================================================================
# 5. PRODUCT / WHAT THEY GET
# =====================================================================
s = slide()
header(s, "What the owner actually sees",
       "Powerful enough for a marketer, simple enough for a busy contractor.")
feats = [
    ("fa-gauge-high", "Today's Priority cockpit", "One screen tells the owner the single most important thing to do today."),
    ("fa-wand-magic-sparkles", "One-tap optimizations", "Approve AI recommendations — or let it run on autopilot with guardrails."),
    ("fa-ban", "Block junk searches", "Plain-English negative keywords stop wasted spend automatically."),
    ("fa-envelope-open-text", "Weekly plain-English digest", "What changed, what it cost, what's next — no jargon."),
    ("fa-shield-halved", "Goal-based guardrails", "Set a target cost per lead and budget cap; agents respect both."),
    ("fa-globe", "Website + SEO on autopilot", "AI writes, refreshes, and optimizes content to win organic leads."),
]
cw = Inches(3.66); ch = Inches(1.7); gx = Inches(0.28); gy = Inches(0.3)
x0 = Inches(0.9); y0 = Inches(2.55)
for i, (ic, t, d) in enumerate(feats):
    r, c = divmod(i, 3)
    x = x0 + c * (cw + gx)
    y = y0 + r * (ch + gy)
    rect(s, x, y, cw, ch, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x + Inches(0.25), y + Inches(0.25), Inches(0.5), Inches(0.5), PRIMARY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x + Inches(0.95), y + Inches(0.2), cw - Inches(1.1), Inches(0.6),
         [[(t, 15, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + Inches(0.28), y + Inches(0.85), cw - Inches(0.5), Inches(0.8),
         [[(d, 12, SLATE, False)]], line_spacing=1.1)
footer(s, 5)

# =====================================================================
# 6. MARKET
# =====================================================================
s = slide()
header(s, "A large, underserved, recurring market",
       "Local service is the last big category without great software-led marketing.")
tam = [
    ("$30B+", "TAM", "US local service businesses' annual digital ad spend"),
    ("$6B+", "SAM", "Home & field service SMBs actively running paid ads"),
    ("32M+", "Businesses", "Home & field service firms in the US"),
]
cw = Inches(3.66); x0 = Inches(0.9); gx = Inches(0.28)
for i, (big, lab, d) in enumerate(tam):
    x = x0 + i * (cw + gx)
    rect(s, x, Inches(2.7), cw, Inches(2.5), INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x, Inches(3.0), cw, Inches(1.0),
         [[(big, 44, ACCENT, True)]], align=PP_ALIGN.CENTER)
    text(s, x, Inches(4.0), cw, Inches(0.5),
         [[(lab, 16, WHITE, True)]], align=PP_ALIGN.CENTER)
    text(s, x + Inches(0.3), Inches(4.5), cw - Inches(0.6), Inches(0.7),
         [[(d, 12, RGBColor(0xCB,0xD5,0xE1), False)]], align=PP_ALIGN.CENTER, line_spacing=1.1)
text(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.0),
     [[("Why now: ", 15, INK, True),
       ("AI agents can finally deliver agency-quality optimization at software margins — and owners now expect it. Whoever becomes the default command center owns the recurring relationship.", 15, SLATE, False)]],
     line_spacing=1.2)
footer(s, 6)

# =====================================================================
# 7. BUSINESS MODEL
# =====================================================================
s = slide()
header(s, "Simple, recurring SaaS — with room to expand",
       "Self-serve tiers land owners; managed and ad-spend-based plans expand revenue.")
plans = [
    ("Free", "$0", "Audit & grader. Top-of-funnel acquisition.", LIGHT, INK),
    ("Growth", "$49/mo", "Single channel, AI optimization, guardrails.", LIGHT, INK),
    ("Pro", "$199/mo", "All channels, autopilot, weekly digests.", PRIMARY, WHITE),
    ("Managed", "$249/mo+", "Done-for-you with human oversight.", INK, WHITE),
]
cw = Inches(2.78); x0 = Inches(0.9); gx = Inches(0.35)
for i, (nm, price, d, bg, fg) in enumerate(plans):
    x = x0 + i * (cw + gx)
    rect(s, x, Inches(2.6), cw, Inches(3.1), bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    sub = RGBColor(0xCB,0xD5,0xE1) if fg == WHITE else SLATE
    text(s, x + Inches(0.3), Inches(2.85), cw - Inches(0.6), Inches(0.5),
         [[(nm, 18, fg, True)]])
    text(s, x + Inches(0.3), Inches(3.35), cw - Inches(0.6), Inches(0.8),
         [[(price, 30, fg, True)]])
    text(s, x + Inches(0.3), Inches(4.25), cw - Inches(0.6), Inches(1.3),
         [[(d, 13, sub, False)]], line_spacing=1.2)
text(s, Inches(0.9), Inches(6.0), Inches(11.5), Inches(0.9),
     [[("Expansion levers: ", 14, INK, True),
       ("more channels per account, % of ad spend on managed tiers, and add-ons (call tracking, reviews, website CRO) — driving net revenue retention above 100%.", 14, SLATE, False)]],
     line_spacing=1.15)
footer(s, 7)

# =====================================================================
# 8. WHY WE WIN
# =====================================================================
s = slide()
header(s, "Why FieldSprout wins",
       "Defensibility compounds with every account and every channel.")
moats = [
    ("Multi-channel by design", "Competitors optimize one channel. Our strategic umbrella moves budget across channels — the owner's whole marketing, not a point tool."),
    ("Built for the non-marketer", "Plain-English UX and guided setup beat both DIY dashboards and intimidating pro tools for the 32M owners who aren't marketers."),
    ("Data & feedback loop", "Every account teaches the agents what works in each trade and geo — a compounding advantage agencies can't match."),
    ("Software economics vs. agencies", "Agency-quality outcomes at SaaS margins; we win on price, speed, and transparency simultaneously."),
]
cw = Inches(5.7); ch = Inches(1.8); x0 = Inches(0.9); y0 = Inches(2.6)
for i, (t, d) in enumerate(moats):
    r, c = divmod(i, 2)
    x = x0 + c * (cw + Inches(0.3))
    y = y0 + r * (ch + Inches(0.25))
    rect(s, x, y, cw, ch, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, y, Inches(0.12), ch, ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x + Inches(0.35), y + Inches(0.22), cw - Inches(0.6), Inches(0.5),
         [[(t, 16, INK, True)]])
    text(s, x + Inches(0.35), y + Inches(0.75), cw - Inches(0.65), Inches(1.0),
         [[(d, 13, SLATE, False)]], line_spacing=1.15)
footer(s, 8)

# =====================================================================
# 9. ROADMAP
# =====================================================================
s = slide()
header(s, "Where we're headed",
       "From multi-channel ads to the full marketing operating system for local pros.")
phases = [
    ("Now", "Google + Facebook Ads and Website/SEO under one AI command center, with guided onboarding and guardrails.", ACCENT),
    ("Next", "Add channels: Google Business Profile, Local Services Ads, email & reviews. Deeper autopilot.", PRIMARY),
    ("Later", "Full marketing OS: lead-to-booked-job attribution, AI front desk, and benchmarked insights across the network.", PRIMARY_D),
]
x = Inches(1.2)
rect(s, x, Inches(3.4), Inches(10.9), Inches(0.05), MUTED)
for i, (ph, d, col) in enumerate(phases):
    cx = x + Inches(0.3) + i * Inches(3.7)
    rect(s, cx, Inches(3.25), Inches(0.35), Inches(0.35), col, shape=MSO_SHAPE.OVAL)
    text(s, cx - Inches(0.3), Inches(2.7), Inches(3.4), Inches(0.5),
         [[(ph, 18, col, True)]])
    text(s, cx - Inches(0.3), Inches(3.85), Inches(3.4), Inches(2.0),
         [[(d, 13, SLATE, False)]], line_spacing=1.2)
footer(s, 9)

# =====================================================================
# 10. CLOSE / ASK
# =====================================================================
s = slide()
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, 0, Inches(0.28), SH, ACCENT)
text(s, Inches(1.0), Inches(1.4), Inches(11), Inches(1.5),
     [[("Become the command center", 44, WHITE, True)],
      [("for local marketing.", 44, ACCENT, True)]], line_spacing=1.0)
text(s, Inches(1.02), Inches(3.5), Inches(10.8), Inches(1.4),
     [[("FieldSprout puts every marketing channel for 32M local service businesses under one", 18, RGBColor(0xCB,0xD5,0xE1), False)],
      [("AI brain — winning the recurring relationship the agencies can't defend.", 18, RGBColor(0xCB,0xD5,0xE1), False)]],
     line_spacing=1.2)
rect(s, Inches(1.0), Inches(5.3), Inches(4.6), Inches(0.9), ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(1.0), Inches(5.35), Inches(4.6), Inches(0.8),
     [[("Let's talk  ·  fieldsprout.io", 18, INK, True)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(1.02), Inches(6.6), Inches(10), Inches(0.5),
     [[("hello@fieldsprout.io", 14, MUTED, False)]])

# ---- save ----
out = sys.argv[1] if len(sys.argv) > 1 else "FieldSprout-Investor-Deck.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
